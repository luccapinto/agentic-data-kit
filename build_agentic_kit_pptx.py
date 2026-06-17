#!/usr/bin/env python3
"""
build_agentic_kit_pptx.py — Apresentação do Agentic Data Kit v4.1.0
"""
from __future__ import annotations

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

BRAND = {
    "ink":    "1B2733",
    "accent": "14517B",
    "accent2":"0E7C66",
    "purple": "5B4A8C",
    "muted":  "6A7787",
    "soft":   "F3F6F9",
    "pos":    "1E7A4D",
    "neg":    "B23A2E",
    "warn":   "9A6B12",
    "bg":     "FFFFFF",
    "font":   "Calibri",
}

def C(key: str) -> RGBColor:
    return RGBColor.from_string(BRAND[key])

SW, SH = Inches(13.333), Inches(7.5)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, x, y, w, h, text, size=18, color="ink", bold=False,
            align=PP_ALIGN.LEFT, font=None, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font or BRAND["font"]
    run.font.color.rgb = C(color)
    return box


def rect(slide, x, y, w, h, color="accent", line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(color)
    if line:
        shape.line.color.rgb = C(color)
    else:
        shape.line.fill.background()
    return shape


def rounded_rect(slide, x, y, w, h, color="soft", line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(color)
    if line_color:
        shape.line.color.rgb = C(line_color)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def header(slide, kicker, title, title_size=26):
    rect(slide, Inches(0.6), Inches(0.55), Inches(12.7), Inches(0.06), "accent")
    textbox(slide, Inches(0.6), Inches(0.68), Inches(11), Inches(0.4),
            kicker.upper(), size=10, color="accent", bold=True)
    textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(1.0),
            title, size=title_size, color="ink", bold=True)


def kpi_card(slide, x, y, w, h, value, label, val_color="accent"):
    card = rounded_rect(slide, x, y, w, h, "soft")
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = C(val_color)
    r.font.name = BRAND["font"]
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(10)
    r2.font.color.rgb = C("muted")
    r2.font.name = BRAND["font"]


def source_line(slide, text):
    textbox(slide, Inches(0.6), Inches(7.05), Inches(12.5), Inches(0.35),
            text, size=9, color="muted", italic=True)


# ── Slides ──────────────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank(prs)
    # Background accent block on left
    rect(s, Inches(0), Inches(0), Inches(4.5), SH, "accent")
    # White title area
    rect(s, Inches(4.5), Inches(0), SW - Inches(4.5), SH, "bg")

    # Left side: logo / version text
    textbox(s, Inches(0.3), Inches(0.4), Inches(3.8), Inches(0.6),
            "AGENTIC DATA KIT", size=14, color="bg", bold=True)
    textbox(s, Inches(0.3), Inches(0.9), Inches(3.8), Inches(0.4),
            "v4.1.0", size=11, color="bg")

    # Left side: tagline
    textbox(s, Inches(0.3), Inches(2.5), Inches(3.8), Inches(2.0),
            "Drop one folder.\nGet a team of AI\ndata specialists.", size=22, color="bg", bold=True)

    textbox(s, Inches(0.3), Inches(5.5), Inches(3.8), Inches(0.6),
            "Claude Code · GitHub Copilot\nOpenCode · Cursor", size=11, color="bg")

    textbox(s, Inches(0.3), Inches(6.8), Inches(3.8), Inches(0.5),
            "Junho 2026", size=10, color="bg")

    # Right side: main title
    textbox(s, Inches(5.0), Inches(2.0), Inches(7.8), Inches(1.8),
            "Uma equipa de IA\npara dados,\nsem complexidade.", size=36, color="ink", bold=True)
    textbox(s, Inches(5.0), Inches(4.0), Inches(7.5), Inches(1.2),
            "5 agentes especialistas · 8 skills on-demand\n3 workflows · 4 plataformas de IA",
            size=16, color="muted")


def slide_problem(prs):
    s = blank(prs)
    header(s, "O problema", "LLMs genéricos não conhecem o teu stack")

    problems = [
        ("Contexto perdido", "Cada sessão começa do zero — sem padrões, sem stack, sem regras da equipa."),
        ("Routing confuso", "O modelo não sabe se é data eng, analytics ou Power BI — tenta tudo de vez."),
        ("Duplicação entre tools", "Copilot, Claude, Cursor — cada um tem instruções diferentes e desalinhadas."),
        ("PII e qualidade ignorados", "Sem guardrails, o modelo sugere pipelines sem masking nem idempotência."),
    ]

    for i, (title, desc) in enumerate(problems):
        row = i // 2
        col = i % 2
        x = Inches(0.6 + col * 6.4)
        y = Inches(2.1 + row * 1.9)
        card = rounded_rect(s, x, y, Inches(6.0), Inches(1.6), "soft")
        rect(s, x, y, Inches(0.18), Inches(1.6), "neg")
        textbox(s, x + Inches(0.3), y + Inches(0.15), Inches(5.5), Inches(0.5),
                title, size=15, bold=True, color="ink")
        textbox(s, x + Inches(0.3), y + Inches(0.55), Inches(5.5), Inches(0.85),
                desc, size=12, color="muted")


def slide_solution(prs):
    s = blank(prs)
    header(s, "A solução", "Contexto especializado, instalado em segundos")

    # Big quote / install command
    cmd_box = rounded_rect(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.0), "ink")
    textbox(s, Inches(0.9), Inches(2.18), Inches(11.5), Inches(0.65),
            "npx @luccapinto/agentic-data-kit@latest init",
            size=18, color="bg", bold=True)

    pillars = [
        ("Agentes por domínio", "Cada especialidade tem o seu próprio agente — sem sobreposição."),
        ("Skills on-demand", "Carregadas só quando necessário — não poluem o contexto."),
        ("Multi-plataforma", "Uma source of truth (.agent/) compila para 4 ferramentas."),
        ("Guardrails de dados", "Idempotência, Write-Audit-Publish e masking de PII embutidos."),
    ]

    for i, (title, desc) in enumerate(pillars):
        x = Inches(0.6 + (i % 2) * 6.4)
        y = Inches(3.4 + (i // 2) * 1.7)
        rect(s, x, y, Inches(0.18), Inches(1.4), "accent2")
        textbox(s, x + Inches(0.3), y + Inches(0.1), Inches(5.8), Inches(0.5),
                title, size=15, bold=True, color="ink")
        textbox(s, x + Inches(0.3), y + Inches(0.55), Inches(5.8), Inches(0.75),
                desc, size=12, color="muted")


def slide_agents(prs):
    s = blank(prs)
    header(s, "Agentes", "5 especialistas com domínios não sobrepostos")

    agents = [
        ("data-engineer", "ETL/ELT · Medallion Architecture · idempotência · Write-Audit-Publish", "accent"),
        ("analytics-engineer", "dbt · Star Schema · DW design · SQL transformations", "accent"),
        ("data-scientist", "Análise · KPIs · ML · forecasting · A/B testing", "accent2"),
        ("powerbi-developer", "TMDL · PBIR · DAX · BPA via Tabular Editor 2", "purple"),
        ("presentation-designer", "reveal.js · python-pptx · sites interativos · brand integration", "accent2"),
    ]

    for i, (name, desc, color) in enumerate(agents):
        y = Inches(2.0 + i * 1.0)
        rect(s, Inches(0.6), y + Inches(0.05), Inches(0.18), Inches(0.75), color)
        num_box = rounded_rect(s, Inches(0.85), y, Inches(0.65), Inches(0.85), "soft")
        textbox(s, Inches(0.87), y + Inches(0.08), Inches(0.58), Inches(0.65),
                f"0{i+1}", size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
        textbox(s, Inches(1.65), y + Inches(0.08), Inches(4.2), Inches(0.38),
                f"@{name}", size=14, bold=True, color="ink")
        textbox(s, Inches(1.65), y + Inches(0.44), Inches(11.0), Inches(0.38),
                desc, size=11, color="muted")


def slide_skills(prs):
    s = blank(prs)
    header(s, "Skills", "8 módulos de conhecimento ativados on-demand")

    skills = [
        ("pbi-semantic-layer-tmdl", "Modelos semânticos Power BI como código (TMDL)"),
        ("pbi-report-layer-pbir", "Reports Power BI como código (PBIR JSON)"),
        ("pbi-quality-rules", "BPA via Tabular Editor 2 CLI — validação automática"),
        ("building-html-presentations", "reveal.js, flex-deck, sites interativos, PDFs"),
        ("generating-pptx", "PowerPoint nativo editável via python-pptx"),
        ("applying-visual-identity", "Brand tokens de DESIGN.md em qualquer output"),
        ("documentation-templates", "Runbooks, dicionários, ADRs, catálogos Power BI"),
        ("creating-agents-and-skills", "Governance para criar e sincronizar agentes/skills"),
    ]

    cols = 2
    col_w = Inches(6.2)
    for i, (name, desc) in enumerate(skills):
        col = i % cols
        row = i // cols
        x = Inches(0.6 + col * 6.5)
        y = Inches(2.1 + row * 1.15)
        rounded_rect(s, x, y, col_w, Inches(1.0), "soft")
        rect(s, x, y, Inches(0.14), Inches(1.0), "accent2")
        textbox(s, x + Inches(0.25), y + Inches(0.08), col_w - Inches(0.35), Inches(0.42),
                name, size=12, bold=True, color="accent")
        textbox(s, x + Inches(0.25), y + Inches(0.48), col_w - Inches(0.35), Inches(0.44),
                desc, size=11, color="muted")


def slide_architecture(prs):
    s = blank(prs)
    header(s, "Arquitetura", "Uma source of truth, quatro plataformas")

    # Center: .agent/ box
    center_x, center_y = Inches(5.4), Inches(2.8)
    center_w, center_h = Inches(2.5), Inches(1.6)
    core = rounded_rect(s, center_x, center_y, center_w, center_h, "accent")
    tf = core.text_frame
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ".agent/"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = C("bg")
    r.font.name = BRAND["font"]
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Source of Truth"
    r2.font.size = Pt(11)
    r2.font.color.rgb = C("bg")
    r2.font.name = BRAND["font"]

    platforms = [
        (Inches(0.6),  Inches(2.9), ".claude/", "Claude Code", "accent2"),
        (Inches(10.3), Inches(2.9), ".github/", "GitHub Copilot", "accent2"),
        (Inches(0.6),  Inches(5.0), ".opencode/", "OpenCode", "purple"),
        (Inches(10.3), Inches(5.0), "AGENTS.md", "Cursor / outros", "purple"),
    ]

    for px, py, folder, label, color in platforms:
        box = rounded_rect(s, px, py, Inches(2.4), Inches(1.3), "soft", line_color=color)
        textbox(s, px + Inches(0.15), py + Inches(0.1), Inches(2.1), Inches(0.45),
                folder, size=14, bold=True, color=color)
        textbox(s, px + Inches(0.15), py + Inches(0.52), Inches(2.1), Inches(0.5),
                label, size=11, color="muted")

    textbox(s, Inches(4.5), Inches(6.4), Inches(4.5), Inches(0.5),
            "scripts/sync_agents.py compila automaticamente",
            size=11, color="muted", align=PP_ALIGN.CENTER, italic=True)


def slide_stats(prs):
    s = blank(prs)
    header(s, "Resumo", "Kit completo, lean e pronto a usar")

    kpis = [
        ("5", "Agentes especialistas", "accent"),
        ("8", "Skills on-demand", "accent2"),
        ("3", "Workflows / slash commands", "purple"),
        ("4", "Plataformas suportadas", "accent"),
    ]

    gap = Inches(0.25)
    margin = Inches(0.6)
    n = len(kpis)
    total_w = SW - margin * 2 - gap * (n - 1)
    w = Emu(int(total_w / n))
    for i, (val, lbl, color) in enumerate(kpis):
        x = margin + Emu(int((w + gap) * i))
        kpi_card(s, x, Inches(2.2), w, Inches(1.8), val, lbl, val_color=color)

    # Chart: lines of context vs value
    cd = CategoryChartData()
    cd.categories = ["data-engineer", "analytics-engineer", "data-scientist", "powerbi-developer", "presentation-designer"]
    cd.add_series("Linhas de instrução", [49, 43, 47, 38, 40])
    gf = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.4), cd
    )
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = False
    for plot in chart.plots:
        for series in plot.series:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = C("accent")

    textbox(s, Inches(0.6), Inches(4.1), Inches(9), Inches(0.3),
            "LINHAS DE INSTRUÇÃO POR AGENTE — lean por design",
            size=10, color="muted", bold=True)
    source_line(s, "Fonte: repositório luccapinto/agentic-data-kit v4.1.0 · MIT License")


def slide_closing(prs):
    s = blank(prs)
    rect(s, Inches(0), Inches(0), SW, SH, "accent")

    textbox(s, Inches(1.0), Inches(1.5), Inches(11), Inches(1.0),
            "AGENTIC DATA KIT", size=14, color="bg", bold=True)

    textbox(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.5),
            "Menos contexto genérico.\nMais especialização.\nMesmos resultados em qualquer ferramenta.",
            size=32, color="bg", bold=True)

    textbox(s, Inches(1.0), Inches(5.3), Inches(11), Inches(0.6),
            "npx @luccapinto/agentic-data-kit@latest init", size=18, color="bg", bold=False)

    textbox(s, Inches(1.0), Inches(6.5), Inches(11), Inches(0.6),
            "github.com/luccapinto/agentic-data-kit · MIT License · v4.1.0",
            size=11, color="bg")


def build(out: str) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_agents(prs)
    slide_skills(prs)
    slide_architecture(prs)
    slide_stats(prs)
    slide_closing(prs)

    prs.save(out)
    print(f"Wrote {out} — {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="agentic_data_kit.pptx")
    args = ap.parse_args()
    build(args.out)
