#!/usr/bin/env python3
"""
build_pptx.py — generate a NATIVE, EDITABLE PowerPoint with python-pptx.

Every element is a real PowerPoint object: editable text, rounded-rectangle KPI cards,
native charts (clustered column, line, doughnut) backed by embedded data, and a table.
This is a working starter — adapt DATA / BRAND / the build() slide calls.

    pip install python-pptx
    python build_pptx.py --out deck.pptx

Demo content is a generic SALES review; replace with real, sourced numbers — never invent data.
"""
from __future__ import annotations

import argparse

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK

# ── BRAND — map these from your DESIGN.md ──────────────────────────────────
BRAND = {
    "ink":    "1B2733",
    "accent": "14517B",
    "accent2":"0E7C66",
    "muted":  "6A7787",
    "soft":   "F3F6F9",
    "pos":    "1E7A4D",
    "neg":    "B23A2E",
    "warn":   "9A6B12",
    "bg":     "FFFFFF",
    "font":   "Inter",
}
def C(key: str) -> RGBColor:
    return RGBColor.from_string(BRAND[key])

# ── DATA — replace with your own ───────────────────────────────────────────
DATA = {
    "title":    "Revenue Performance",
    "subtitle": "An exploratory analysis of FY2026 Q1 bookings",
    "date":     "June 2026 · Confidential",
    "kpis":     [("$4.2M", "Bookings"), ("+12%", "vs. plan"), ("28%", "Win rate"), ("$9.1M", "Pipeline")],
    "quarters": ["Q2/25", "Q3/25", "Q4/25", "Q1/26"],
    "trend":    [3.1, 3.5, 3.8, 4.2],
    "regions":  ["North America", "EMEA", "LATAM", "APAC"],
    "bookings": [1.8, 1.1, 0.7, 0.6],
    "segments": ["Enterprise", "Mid-market", "SMB"],
    "attain_plan":   [122, 99, 73],
    "attain_actual": [118, 104, 88],
    "table": [
        ["Region", "Bookings", "Attainment", "Deals", "Win rate"],
        ["North America", "$1.8M", "118%", "402", "31%"],
        ["EMEA", "$1.1M", "104%", "318", "29%"],
        ["LATAM", "$0.7M", "96%", "266", "26%"],
        ["APAC", "$0.6M", "88%", "254", "22%"],
        ["Total", "$4.2M", "112%", "1,240", "28%"],
    ],
}

SW, SH = Inches(13.333), Inches(7.5)  # 16:9


# ── Helpers ────────────────────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, x, y, w, h, text, size=18, color="ink", bold=False,
            align=PP_ALIGN.LEFT, font=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font or BRAND["font"]
    run.font.color.rgb = C(color)
    return box


def accent_bar(slide, x, y, w, h, color="accent"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid(); bar.fill.fore_color.rgb = C(color)
    bar.line.fill.background()
    return bar


def kpi_card(slide, x, y, w, h, value, label):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = C("soft")
    card.line.color.rgb = C("soft"); card.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.16)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = value
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = C("accent"); r.font.name = BRAND["font"]
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = label.upper()
    r2.font.size = Pt(10); r2.font.color.rgb = C("muted"); r2.font.name = BRAND["font"]
    return card


def header(slide, kicker, title):
    accent_bar(slide, Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.09))
    textbox(slide, Inches(0.6), Inches(0.62), Inches(11), Inches(0.4),
            kicker.upper(), size=11, color="accent", bold=True)
    textbox(slide, Inches(0.6), Inches(0.95), Inches(12.1), Inches(1.0),
            title, size=24, color="ink", bold=True)


def style_chart(chart, colors, legend=True):
    chart.has_title = False
    if legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10); chart.legend.font.color.rgb = C("muted")
    else:
        chart.has_legend = False
    for i, plot in enumerate(chart.plots):
        plot.has_data_labels = False
        for j, series in enumerate(plot.series):
            try:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = colors[j % len(colors)]
            except Exception:
                pass


def source_line(slide, text):
    textbox(slide, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
            text, size=9, color="muted")


# ── Slides ───────────────────────────────────────────────────────────────
def slide_title(prs):
    s = blank(prs)
    accent_bar(s, Inches(0.6), Inches(2.2), Inches(2.0), Inches(0.12))
    textbox(s, Inches(0.6), Inches(2.5), Inches(11), Inches(1.3), DATA["title"], size=48, bold=True)
    textbox(s, Inches(0.62), Inches(3.7), Inches(10), Inches(0.8), DATA["subtitle"], size=20, color="muted")
    textbox(s, Inches(0.62), Inches(6.4), Inches(10), Inches(0.5), DATA["date"], size=12, color="muted")


def slide_agenda(prs):
    s = blank(prs)
    header(s, "Agenda", "What this report covers")
    items = ["Executive summary & KPIs", "Bookings trend", "Concentration by region & segment",
             "Region-level detail", "Recommendations"]
    for i, it in enumerate(items):
        y = Inches(2.0 + i * 0.85)
        textbox(s, Inches(0.7), y, Inches(0.8), Inches(0.7), f"{i+1:02d}", size=26, color="accent", bold=True)
        textbox(s, Inches(1.7), y + Inches(0.08), Inches(10), Inches(0.6), it, size=18)


def slide_kpis(prs):
    s = blank(prs)
    header(s, "Executive summary", "Q1 bookings beat plan by 12%, led by enterprise")
    n = len(DATA["kpis"]); gap = Inches(0.3); margin = Inches(0.6)
    total_w = SW - margin * 2 - gap * (n - 1)
    w = Emu(int(total_w / n)); h = Inches(1.6)
    for i, (val, lbl) in enumerate(DATA["kpis"]):
        x = margin + Emu(int((w + gap) * i))
        kpi_card(s, x, Inches(2.2), w, h, val, lbl)
    textbox(s, Inches(0.6), Inches(4.3), Inches(12), Inches(1.5),
            "Upside is real but narrow: two regions and the enterprise segment carry most of it, "
            "while SMB attainment (73%) and discounting are early warning signs for H2.",
            size=15, color="muted")
    source_line(s, "Source: company data (gold.fct_bookings) — figures illustrative")


def slide_trend(prs):
    s = blank(prs)
    header(s, "Trend", "Bookings grew every quarter, up 35% since Q2")
    cd = CategoryChartData()
    cd.categories = DATA["quarters"]
    cd.add_series("Bookings ($M)", DATA["trend"])
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.8), Inches(2.0),
                            Inches(11.7), Inches(4.6), cd)
    style_chart(gf.chart, [C("accent")], legend=False)
    source_line(s, "Exhibit 1 · Source: company data; n = 1,240 closed-won deals")


def slide_region_segment(prs):
    s = blank(prs)
    header(s, "Concentration", "Two regions drive ~70%; enterprise leads, SMB lags")
    # Exhibit 2 — clustered column: bookings by region
    cd = CategoryChartData(); cd.categories = DATA["regions"]
    cd.add_series("Bookings ($M)", DATA["bookings"])
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(2.0),
                            Inches(6.0), Inches(4.4), cd)
    style_chart(gf.chart, [C("accent")], legend=False)
    textbox(s, Inches(0.8), Inches(1.75), Inches(5), Inches(0.3), "EXHIBIT 2 · Bookings by region",
            size=10, color="muted", bold=True)
    # Exhibit 3 — clustered column: plan vs actual attainment by segment
    cd2 = CategoryChartData(); cd2.categories = DATA["segments"]
    cd2.add_series("Plan", DATA["attain_plan"])
    cd2.add_series("Actual", DATA["attain_actual"])
    gf2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(2.0),
                             Inches(5.6), Inches(4.4), cd2)
    style_chart(gf2.chart, [C("muted"), C("accent2")], legend=True)
    textbox(s, Inches(7.1), Inches(1.75), Inches(5), Inches(0.3), "EXHIBIT 3 · Attainment by segment (%)",
            size=10, color="muted", bold=True)
    source_line(s, "Source: company data; target set at FY2026 plan")


def slide_pie(prs):
    s = blank(prs)
    header(s, "Mix", "North America and EMEA account for most bookings")
    cd = CategoryChartData(); cd.categories = DATA["regions"]
    cd.add_series("Share", DATA["bookings"])
    gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(3.3), Inches(2.0),
                            Inches(6.7), Inches(4.6), cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False; chart.legend.font.size = Pt(11)
    points = chart.plots[0].series[0].points
    palette = [C("accent"), C("accent2"), C("warn"), C("muted")]
    for i, pt in enumerate(points):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = palette[i % len(palette)]
    source_line(s, "Exhibit 4 · Source: company data; shares of $4.2M total")


def slide_table(prs):
    s = blank(prs)
    header(s, "Appendix", "Region-level detail")
    rows, cols = len(DATA["table"]), len(DATA["table"][0])
    tbl = s.shapes.add_table(rows, cols, Inches(0.7), Inches(2.0),
                             Inches(11.9), Inches(0.5 * rows)).table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = DATA["table"][r][c]
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            run = para.runs[0]
            run.font.size = Pt(12); run.font.name = BRAND["font"]
            is_edge = (r == 0 or r == rows - 1)
            run.font.bold = is_edge
            run.font.color.rgb = C("bg") if r == 0 else C("ink")
            cell.fill.solid()
            cell.fill.fore_color.rgb = C("accent") if r == 0 else (C("soft") if r == rows - 1 else C("bg"))
    source_line(s, "Exhibit 5 · Source: company data — replace with real, sourced values")


def slide_recs(prs):
    s = blank(prs)
    header(s, "Recommendations", "Four moves to protect the upside")
    recs = [
        ("Scale the enterprise motion", "Reallocate capacity to the segment that beat plan."),
        ("Broaden mid-market pipeline", "Reduce concentration beyond the top two regions."),
        ("Fix the SMB funnel", "Tighten qualification and discount guardrails."),
        ("Re-forecast H2 with scenarios", "Model downside on SMB and discounting; set checkpoints."),
    ]
    for i, (h, p) in enumerate(recs):
        y = Inches(2.1 + i * 1.1)
        textbox(s, Inches(0.7), y, Inches(0.8), Inches(0.7), f"{i+1}", size=28, color="accent", bold=True)
        textbox(s, Inches(1.6), y, Inches(11), Inches(0.5), h, size=17, bold=True)
        textbox(s, Inches(1.6), y + Inches(0.45), Inches(11), Inches(0.5), p, size=13, color="muted")


def build(out: str) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    slide_title(prs)
    slide_agenda(prs)
    slide_kpis(prs)
    slide_trend(prs)
    slide_region_segment(prs)
    slide_pie(prs)
    slide_table(prs)
    slide_recs(prs)
    prs.save(out)
    print(f"Wrote {out} · {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="deck.pptx")
    build(ap.parse_args().out)
